import streamlit as st
from azure.search.documents import SearchClient
from azure.core.credentials import AzureKeyCredential
from openai import AzureOpenAI
from dotenv import load_dotenv
import os
from io import BytesIO
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
from docx import Document
from docx.shared import Pt, RGBColor, Inches
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Inches, Pt as PptPt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor as PptRGBColor

# 환경 변수 로드
load_dotenv()

# Azure 설정
SEARCH_ENDPOINT = os.getenv("AZURE_SEARCH_ENDPOINT")
SEARCH_KEY = os.getenv("AZURE_SEARCH_KEY")
SEARCH_INDEX = os.getenv("AZURE_SEARCH_INDEX")
OPENAI_API_KEY = os.getenv("AZURE_OPENAI_KEY")
OPENAI_DEPLOYMENT = os.getenv("AZURE_OPENAI_DEPLOYMENT")
OPENAI_ENDPOINT = os.getenv("AZURE_OPENAI_ENDPOINT")

# 클라이언트 초기화
@st.cache_resource
def get_search_client():
    return SearchClient(
        endpoint=SEARCH_ENDPOINT,
        index_name=SEARCH_INDEX,
        credential=AzureKeyCredential(SEARCH_KEY)
    )

@st.cache_resource
def get_openai_client():
    return AzureOpenAI(
        api_key=OPENAI_API_KEY,
        api_version="2024-12-01-preview",
        azure_endpoint=OPENAI_ENDPOINT
    )

search_client = get_search_client()
openai_client = get_openai_client()

# GPT 호출 함수
def ask_openai(prompt=None, messages=None):
    try:
        if messages:
            response = openai_client.chat.completions.create(
                model=OPENAI_DEPLOYMENT,
                messages=messages,
                temperature=0.7,
                max_tokens=2000
            )
        elif prompt:
            response = openai_client.chat.completions.create(
                model=OPENAI_DEPLOYMENT,
                messages=[
                    {"role": "system", "content": "당신은 산업 분석 전문가입니다."},
                    {"role": "user", "content": prompt}
                ],
                temperature=0.7,
                max_tokens=2000
            )
        else:
            raise ValueError("prompt 또는 messages 중 하나는 반드시 필요합니다.")
        return response.choices[0].message.content
    except Exception as e:
        st.error(f"❌ GPT 호출 중 오류 발생: {str(e)}")
        return None

# PDF 생성 함수
def create_pdf(content, title):
    buffer = BytesIO()
    doc = SimpleDocTemplate(buffer, pagesize=A4, rightMargin=72, leftMargin=72, topMargin=72, bottomMargin=18)
    
    # 한글 폰트 등록
    try:
        pdfmetrics.registerFont(TTFont('NanumGothic', 'NanumGothic-Regular.ttf'))
        pdfmetrics.registerFont(TTFont('NanumGothic-Bold', 'NanumGothic-Bold.ttf'))
        pdfmetrics.registerFont(TTFont('NanumGothic-ExtraBold', 'NanumGothic-ExtraBold.ttf'))
        font_normal = 'NanumGothic'
        font_bold = 'NanumGothic-Bold'
        font_extra_bold = 'NanumGothic-ExtraBold'
    except Exception as e:
        st.warning(f"⚠️ 한글 폰트 로딩 실패: {str(e)}. 기본 폰트를 사용합니다.")
        font_normal = 'Helvetica'
        font_bold = 'Helvetica-Bold'
        font_extra_bold = 'Helvetica-Bold'
    
    styles = getSampleStyleSheet()
    
    # 커스텀 스타일 정의
    styles.add(ParagraphStyle(
        name='CustomTitle', 
        parent=styles['Heading1'], 
        fontSize=24, 
        textColor='blue', 
        spaceAfter=30, 
        alignment=1,
        fontName=font_extra_bold,
        leading=30
    ))
    
    styles.add(ParagraphStyle(
        name='CustomHeading', 
        parent=styles['Heading2'], 
        fontSize=16, 
        spaceAfter=12, 
        spaceBefore=20,
        fontName=font_bold,
        leading=20
    ))
    
    styles.add(ParagraphStyle(
        name='CustomBody', 
        parent=styles['BodyText'], 
        fontSize=11, 
        spaceAfter=12,
        fontName=font_normal,
        leading=18,
        wordWrap='CJK'
    ))
    
    story = []
    story.append(Paragraph(title, styles['CustomTitle']))
    story.append(Spacer(1, 0.5*inch))
    
    lines = content.split('\n')
    for line in lines:
        line = line.strip()
        if not line:
            continue
            
        if line.startswith('## '):
            story.append(Spacer(1, 0.3*inch))
            story.append(Paragraph(line.replace('## ', ''), styles['CustomHeading']))
        elif line.startswith('# '):
            continue
        else:
            story.append(Paragraph(line, styles['CustomBody']))
    
    try:
        doc.build(story)
    except Exception as e:
        st.error(f"PDF 생성 중 오류: {str(e)}")
        return None
    
    buffer.seek(0)
    return buffer

# Word 생성 함수
def create_word(content, title):
    doc = Document()
    
    title_para = doc.add_heading(title, level=0)
    title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    title_run = title_para.runs[0]
    title_run.font.size = Pt(24)
    title_run.font.color.rgb = RGBColor(0, 112, 192)
    
    doc.add_paragraph()
    
    lines = content.split('\n')
    for line in lines:
        line = line.strip()
        if not line:
            continue
            
        if line.startswith('## '):
            doc.add_heading(line.replace('## ', ''), level=1)
        elif line.startswith('# '):
            continue
        else:
            para = doc.add_paragraph(line)
            para.style.font.size = Pt(11)
    
    buffer = BytesIO()
    doc.save(buffer)
    buffer.seek(0)
    return buffer

# PowerPoint 생성 함수
def create_ppt(content, title):
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = title
    
    lines = content.split('\n')
    current_slide = None
    current_content = []
    
    for line in lines:
        line = line.strip()
        if not line:
            continue
            
        if line.startswith('## '):
            if current_slide and current_content:
                text_frame = current_slide.placeholders[1].text_frame
                for content_line in current_content:
                    p = text_frame.add_paragraph()
                    p.text = content_line
                    p.level = 0
                    p.font.size = PptPt(14)
            
            bullet_slide_layout = prs.slide_layouts[1]
            current_slide = prs.slides.add_slide(bullet_slide_layout)
            current_slide.shapes.title.text = line.replace('## ', '')
            current_content = []
        elif line.startswith('# '):
            continue
        elif current_slide:
            current_content.append(line)
    
    if current_slide and current_content:
        text_frame = current_slide.placeholders[1].text_frame
        for content_line in current_content:
            p = text_frame.add_paragraph()
            p.text = content_line
            p.level = 0
            p.font.size = PptPt(14)
    
    buffer = BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer

# Streamlit UI 시작
st.set_page_config(page_title="RAG 기반 산업군 추천 및 분석", layout="wide")

# 상태 초기화
def initialize_session_state():
    defaults = {
        "keyword": "",
        "search_results": [],
        "recommendations_raw": "",
        "recommendation_list": [],
        "selected_industry": "",
        "chat_history": [],
        "report_sections": [],  # 섹션별 작성 상태 추적
        "report_final": "",
        "report_completed": False,
        "generated_file": None,
        "current_section_index": 0,  # 현재 작성 중인 섹션 인덱스
    }
    for key, default in defaults.items():
        if key not in st.session_state:
            st.session_state[key] = default

initialize_session_state()

# 보고서 섹션 정의
REPORT_SECTIONS = [
    {"title": "산업 개요 및 시장 동향", "description": "산업의 정의, 현황, 최신 트렌드"},
    {"title": "시장 규모 및 성장 전망", "description": "구체적인 시장 규모 데이터와 성장 예측"},
    {"title": "주요 경쟁사 분석", "description": "주요 플레이어들의 현황과 포지셔닝"},
    {"title": "핵심 기술 및 혁신 동향", "description": "기술적 발전과 혁신 사례"},
    {"title": "타겟 고객 및 시장 세그먼트", "description": "주요 고객층과 시장 세분화"},
    {"title": "사업 기회 및 진입 전략", "description": "시장 기회와 전략적 접근"},
    {"title": "리스크 요인 및 대응 방안", "description": "잠재적 위험과 완화 전략"},
    {"title": "결론 및 향후 전망", "description": "종합 분석과 미래 전망"}
]

# 초기화 버튼
with st.sidebar:
    st.header("⚙️ 설정")
    if st.button("🔄 전체 초기화", use_container_width=True):
        for key in list(st.session_state.keys()):
            del st.session_state[key]
        st.rerun()
    
    # 보고서 작성 진행 상황 표시
    if st.session_state.report_sections:
        st.markdown("---")
        st.subheader("📊 보고서 작성 진행")
        progress = len(st.session_state.report_sections) / len(REPORT_SECTIONS)
        st.progress(progress)
        st.metric("작성 완료", f"{len(st.session_state.report_sections)}/{len(REPORT_SECTIONS)} 섹션")
        
        st.markdown("**완료된 섹션:**")
        for i, section in enumerate(st.session_state.report_sections):
            st.markdown(f"✅ {i+1}. {section['title']}")

# 보고서 완료 후 화면
if st.session_state.report_completed:
    st.markdown("""
    <style>
    .report-title {
        font-size: 2.5rem;
        font-weight: bold;
        text-align: center;
        margin: 2rem 0;
        color: #1f77b4;
        border-bottom: 3px solid #1f77b4;
        padding-bottom: 1rem;
    }
    .report-content {
        font-size: 1.1rem;
        line-height: 1.8;
    }
    </style>
    """, unsafe_allow_html=True)
    
    report_title = f"{st.session_state.selected_industry} 시장 분석 및 사업 제안 보고서"
    
    st.markdown(f'<h1 class="report-title">{report_title}</h1>', unsafe_allow_html=True)
    st.markdown("---")
    st.markdown('<div class="report-content">', unsafe_allow_html=True)
    st.markdown(st.session_state.report_final)
    st.markdown('</div>', unsafe_allow_html=True)
    st.markdown("---")
    
    st.subheader("📁 보고서 파일 생성")
    
    col1, col2 = st.columns([1, 2])
    with col1:
        file_format = st.selectbox(
            "파일 형식 선택",
            ["PDF", "Word", "PowerPoint"],
            key="file_format_selector"
        )
    
    with col2:
        if st.button("📄 보고서 파일 생성", use_container_width=True, type="primary", key="generate_file_btn"):
            with st.spinner(f"🔄 {file_format} 파일 생성 중..."):
                try:
                    if file_format == "PDF":
                        file_buffer = create_pdf(st.session_state.report_final, report_title)
                        mime_type = "application/pdf"
                        file_ext = "pdf"
                    elif file_format == "Word":
                        file_buffer = create_word(st.session_state.report_final, report_title)
                        mime_type = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
                        file_ext = "docx"
                    elif file_format == "PowerPoint":
                        file_buffer = create_ppt(st.session_state.report_final, report_title)
                        mime_type = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
                        file_ext = "pptx"
                    
                    if file_buffer:
                        st.session_state.generated_file = {
                            "buffer": file_buffer,
                            "mime": mime_type,
                            "ext": file_ext,
                            "format": file_format
                        }
                        st.success(f"✅ {file_format} 파일 생성 완료!")
                except Exception as e:
                    st.error(f"❌ 파일 생성 중 오류 발생: {str(e)}")
    
    if st.session_state.generated_file:
        st.markdown("---")
        st.info(f"📥 **{st.session_state.generated_file['format']}** 파일이 준비되었습니다.")
        
        col1, col2, col3 = st.columns([2, 2, 1])
        with col1:
            st.download_button(
                label=f"💾 {st.session_state.generated_file['format']} 파일 다운로드",
                data=st.session_state.generated_file["buffer"],
                file_name=f"보고서_{st.session_state.selected_industry}_{st.session_state.keyword}.{st.session_state.generated_file['ext']}",
                mime=st.session_state.generated_file["mime"],
                use_container_width=True,
                key=f"download_{st.session_state.generated_file['format']}"
            )
        with col2:
            if st.button("🔄 다른 형식으로 생성", use_container_width=True, key="reset_file_btn"):
                st.session_state.generated_file = None
                st.rerun()
    
    st.markdown("---")
    if st.button("🔄 새로운 분석 시작", use_container_width=True, key="new_analysis_btn"):
        for key in list(st.session_state.keys()):
            del st.session_state[key]
        st.rerun()
    
    st.stop()

# 앱 제목
st.title("🔍 관심 키워드 기반 RAG 검색 → 산업군 추천 → GPT 분석")
st.markdown("---")

# 1. 관심 키워드 입력
st.subheader("1️⃣ 관심 키워드 입력")
keyword_input = st.text_input(
    "분석하고 싶은 키워드를 입력하세요",
    st.session_state.keyword,
    placeholder="예: 인공지능, 전기차, 바이오테크 등"
)

if keyword_input != st.session_state.keyword:
    st.session_state.keyword = keyword_input
    st.session_state.search_results = []
    st.session_state.recommendations_raw = ""
    st.session_state.recommendation_list = []
    st.session_state.selected_industry = ""
    st.session_state.chat_history = []

# 2. RAG 검색 수행
if st.session_state.keyword and not st.session_state.search_results:
    with st.spinner("🔍 RAG 검색 중..."):
        try:
            results = search_client.search(st.session_state.keyword, top=5)
            docs = []
            result_list = list(results)
            
            if not result_list:
                st.warning("⚠️ 검색 결과가 없습니다. 다른 키워드를 시도해보세요.")
            else:
                for i, result in enumerate(result_list):
                    title = result.get("title", f"문서 {i+1}")
                    chunk = result.get("chunk", "")
                    docs.append(f"{title}\n{chunk}")
                st.session_state.search_results = docs
        except Exception as e:
            st.error(f"❌ 검색 중 오류 발생: {str(e)}")

# 2-1. RAG 검색 결과 표시
if st.session_state.search_results:
    st.markdown("---")
    st.subheader("2️⃣ RAG 검색 결과")
    st.info(f"📊 총 {len(st.session_state.search_results)}개의 문서가 검색되었습니다.")
    
    for i, doc in enumerate(st.session_state.search_results):
        title = doc.split('\n')[0] if '\n' in doc else f"문서 {i+1}"
        chunk = doc.split('\n', 1)[1] if '\n' in doc else doc
        with st.expander(f"📄 {title}"):
            st.markdown(chunk[:500] + "..." if len(chunk) > 500 else chunk)

# 3. GPT 산업군 추천
if st.session_state.search_results and not st.session_state.recommendations_raw:
    with st.spinner("🤖 GPT가 산업군을 분석 중..."):
        combined_text = "\n\n".join(st.session_state.search_results)
        prompt = f"""
다음은 '{st.session_state.keyword}' 키워드에 대해 검색된 문서 내용입니다:

{combined_text}

이 정보를 바탕으로 관련된 유망 산업군을 5개 추천해 주세요.
각 산업군은 다음 형식으로 작성해 주세요:
- 산업군명: 간단한 설명 (1-2문장)

반드시 각 줄은 "산업군명:"으로 시작하고 그 뒤에 설명이 오도록 작성해 주세요.
        """
        recommendations = ask_openai(prompt=prompt)
        if recommendations:
            st.session_state.recommendations_raw = recommendations
            recommendation_list = []
            for line in recommendations.split("\n"):
                line = line.strip("-• ").strip()
                if ':' in line and any(c.isalpha() for c in line):
                    title = line.split(':')[0].strip()
                    if title:
                        recommendation_list.append(title)
            st.session_state.recommendation_list = recommendation_list[:5]
        st.rerun()

# 4. 추천 결과 출력 및 산업군 선택
if st.session_state.recommendation_list:
    st.markdown("---")
    st.subheader("3️⃣ GPT 추천 산업군")
    st.markdown(st.session_state.recommendations_raw)
    
    st.markdown("---")
    current_index = 0
    if st.session_state.selected_industry and st.session_state.selected_industry in st.session_state.recommendation_list:
        current_index = st.session_state.recommendation_list.index(st.session_state.selected_industry)
    
    selected = st.selectbox(
        "🎯 분석하고 싶은 산업군을 선택하세요",
        st.session_state.recommendation_list,
        index=current_index
    )
    
    if selected != st.session_state.selected_industry:
        st.session_state.selected_industry = selected
        if st.session_state.chat_history:
            if st.checkbox("이전 대화 기록 유지"):
                pass
            else:
                st.session_state.chat_history = []

# 5. GPT 질의응답
if st.session_state.selected_industry:
    st.markdown("---")
    st.subheader(f"4️⃣ '{st.session_state.selected_industry}' 산업군 GPT 질의응답")
    
    user_question = st.text_input(
        "💬 질문을 입력하세요",
        key="chat_input",
        placeholder="예: 이 산업의 주요 트렌드는 무엇인가요?"
    )
    
    col1, col2 = st.columns([1, 5])
    with col1:
        ask_button = st.button("📤 질문하기", use_container_width=True)
    
    if ask_button and user_question.strip():
        with st.spinner("🤖 GPT가 답변 생성 중..."):
            messages = [
                {"role": "system", "content": f"당신은 '{st.session_state.selected_industry}' 산업군의 시장 분석 전문가입니다."}
            ]
            for q, a in st.session_state.chat_history:
                messages.append({"role": "user", "content": q})
                messages.append({"role": "assistant", "content": a})
            messages.append({"role": "user", "content": user_question})
            
            answer = ask_openai(messages=messages)
            if answer:
                st.session_state.chat_history.append((user_question, answer))
                st.rerun()

# 6. 대화 기록 출력
if st.session_state.chat_history:
    st.markdown("---")
    st.subheader("🗂️ 대화 기록 (최신순)")
    for i, (q, a) in enumerate(reversed(st.session_state.chat_history)):
        idx = len(st.session_state.chat_history) - i
        with st.expander(f"Q{idx}: {q[:50]}...", expanded=(i == 0)):
            st.markdown(f"**질문:** {q}")
            st.markdown(f"**답변:** {a}")

# 7. 보고서 작성 흐름 (섹션별 순차 작성)
if st.session_state.chat_history:
    st.markdown("---")
    st.subheader("5️⃣ 종합 보고서 작성")
    
    # 현재 진행 상황 표시
    if st.session_state.report_sections:
        progress = len(st.session_state.report_sections) / len(REPORT_SECTIONS)
        st.progress(progress)
        st.info(f"📝 보고서 작성 진행: {len(st.session_state.report_sections)}/{len(REPORT_SECTIONS)} 섹션 완료")
    
    # 현재 작성할 섹션 정보
    if st.session_state.current_section_index < len(REPORT_SECTIONS):
        current_section = REPORT_SECTIONS[st.session_state.current_section_index]
        
        st.markdown(f"""
        ### 📌 다음 작성 섹션
        **{st.session_state.current_section_index + 1}. {current_section['title']}**
        
        *{current_section['description']}*
        """)
        
        # 보고서 작성 시작 또는 계속
        button_label = "📝 보고서 작성 시작" if not st.session_state.report_sections else f"➕ '{current_section['title']}' 섹션 작성"
        
        if st.button(button_label, use_container_width=True, type="primary"):
            with st.spinner(f"🔄 '{current_section['title']}' 작성 중..."):
                # 메시지 구성
                messages = [
                    {"role": "system", "content": f"당신은 '{st.session_state.selected_industry}' 산업군에 대한 전문적인 시장 분석 보고서를 작성하는 전문가입니다."}
                ]
                
                # 이전 대화 내역 추가
                for q, a in st.session_state.chat_history:
                    messages.append({"role": "user", "content": q})
                    messages.append({"role": "assistant", "content": a})
                
                # 이미 작성된 섹션이 있으면 컨텍스트로 제공
                if st.session_state.report_sections:
                    previous_content = "\n\n".join([
                        f"## {section['title']}\n{section['content']}" 
                        for section in st.session_state.report_sections
                    ])
                    messages.append({
                        "role": "assistant", 
                        "content": f"[지금까지 작성된 보고서 내용]\n\n{previous_content}"
                    })
                
                # 현재 섹션 작성 요청
                section_prompt = f"""
이제 보고서의 다음 섹션을 작성해주세요:

**섹션 {st.session_state.current_section_index + 1}: {current_section['title']}**

이 섹션에서 다룰 내용: {current_section['description']}

**작성 지침:**
1. "## {current_section['title']}" 형식으로 섹션 제목을 시작하세요
2. 이전 섹션들과 자연스럽게 연결되도록 작성하세요
3. 구체적인 데이터, 사례, 분석을 포함하세요
4. 전문적이고 설득력 있게 작성하세요
5. 넘버링이나 리스트를 사용할 때는 일관성을 유지하세요
6. 마무리 멘트 없이 섹션 내용만 작성하세요

지금까지의 대화 내용과 이전 섹션들을 참고하여 '{current_section['title']}' 섹션을 상세히 작성해주세요.
"""
                messages.append({"role": "user", "content": section_prompt})
                
                # GPT 호출
                section_content = ask_openai(messages=messages)
                
                if section_content:
                    # 섹션 저장
                    st.session_state.report_sections.append({
                        "title": current_section['title'],
                        "content": section_content
                    })
                    st.session_state.current_section_index += 1
                    st.success(f"✅ '{current_section['title']}' 섹션 작성 완료!")
                    st.rerun()
    
    # 작성된 보고서 미리보기
    if st.session_state.report_sections:
        st.markdown("---")
        st.subheader("📄 작성 중인 보고서 미리보기")
        
        # 전체 보고서 조합
        full_report = f"# {st.session_state.selected_industry} 시장 분석 및 사업 제안 보고서\n\n"
        full_report += "\n\n".join([
            f"## {section['title']}\n{section['content']}" 
            for section in st.session_state.report_sections
        ])
        
        with st.expander("🔍 현재까지 작성된 전체 보고서 보기", expanded=False):
            st.markdown(full_report)
        
        # 각 섹션별 미리보기
        st.markdown("### 작성된 섹션 목록")
        for i, section in enumerate(st.session_state.report_sections):
            with st.expander(f"✅ {i+1}. {section['title']}", expanded=False):
                st.markdown(section['content'])
        
        # 모든 섹션 작성 완료 시 최종 완료 버튼
        if st.session_state.current_section_index >= len(REPORT_SECTIONS):
            st.markdown("---")
            st.success("🎉 모든 섹션 작성이 완료되었습니다!")
            
            col1, col2 = st.columns(2)
            with col1:
                if st.button("✅ 보고서 최종 완료", use_container_width=True, type="primary"):
                    # 최종 보고서 조합
                    final_report = f"# {st.session_state.selected_industry} 시장 분석 및 사업 제안 보고서\n\n"
                    final_report += "\n\n".join([
                        f"## {section['title']}\n{section['content']}" 
                        for section in st.session_state.report_sections
                    ])
                    st.session_state.report_final = final_report
                    st.session_state.report_completed = True
                    st.rerun()
            
            with col2:
                if st.button("🔄 특정 섹션 수정", use_container_width=True):
                    st.info("💡 수정하고 싶은 섹션의 인덱스를 아래에서 선택하세요.")
                    section_to_edit = st.selectbox(
                        "수정할 섹션 선택",
                        range(len(st.session_state.report_sections)),
                        format_func=lambda x: f"{x+1}. {st.session_state.report_sections[x]['title']}"
                    )
                    
                    if st.button("📝 선택한 섹션 재작성", key="rewrite_section"):
                        with st.spinner(f"🔄 '{st.session_state.report_sections[section_to_edit]['title']}' 재작성 중..."):
                            selected_section = REPORT_SECTIONS[section_to_edit]
                            
                            messages = [
                                {"role": "system", "content": f"당신은 '{st.session_state.selected_industry}' 산업군에 대한 전문적인 시장 분석 보고서를 작성하는 전문가입니다."}
                            ]
                            
                            for q, a in st.session_state.chat_history:
                                messages.append({"role": "user", "content": q})
                                messages.append({"role": "assistant", "content": a})
                            
                            # 이전/이후 섹션 컨텍스트 제공
                            context_sections = []
                            if section_to_edit > 0:
                                context_sections.append(f"[이전 섹션]\n## {st.session_state.report_sections[section_to_edit-1]['title']}\n{st.session_state.report_sections[section_to_edit-1]['content']}")
                            if section_to_edit < len(st.session_state.report_sections) - 1:
                                context_sections.append(f"[다음 섹션]\n## {st.session_state.report_sections[section_to_edit+1]['title']}\n{st.session_state.report_sections[section_to_edit+1]['content']}")
                            
                            if context_sections:
                                messages.append({"role": "assistant", "content": "\n\n".join(context_sections)})
                            
                            rewrite_prompt = f"""
'{selected_section['title']}' 섹션을 다시 작성해주세요.

이전 내용:
{st.session_state.report_sections[section_to_edit]['content']}

**개선 사항:**
- 더 구체적인 데이터와 분석 추가
- 논리적 흐름 강화
- 이전/다음 섹션과의 연결성 개선
- 전문성과 설득력 향상

"## {selected_section['title']}" 형식으로 시작하여 섹션을 재작성해주세요.
"""
                            messages.append({"role": "user", "content": rewrite_prompt})
                            
                            new_content = ask_openai(messages=messages)
                            if new_content:
                                st.session_state.report_sections[section_to_edit]['content'] = new_content
                                st.success(f"✅ '{selected_section['title']}' 섹션이 재작성되었습니다!")
                                st.rerun()